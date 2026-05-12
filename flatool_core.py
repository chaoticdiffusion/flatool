import hashlib
import io
import re
import zipfile
from dataclasses import dataclass
from typing import Iterable

import fitz
from PIL import Image
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.util import Inches


VALID_LICENSE_HASHES = {
    hashlib.sha256("FLATOOL-BETA-2026".encode()).hexdigest(),
    hashlib.sha256("CCC-FOUNDER-2026".encode()).hexdigest(),
}


@dataclass
class MemoryUpload:
    name: str
    data: bytes

    def getvalue(self) -> bytes:
        return self.data


def natural_sort_key(value: str) -> list[object]:
    return [
        int(part) if part.isdigit() else part.lower()
        for part in re.split(r"([0-9]+)", value)
    ]


def is_valid_license(license_key: str) -> bool:
    normalized = license_key.strip().upper()
    return hashlib.sha256(normalized.encode()).hexdigest() in VALID_LICENSE_HASHES


def build_output_name(uploaded_files) -> str:
    first_file = sorted(uploaded_files, key=lambda file: natural_sort_key(file.name))[0]
    return build_result_name(first_file.name, ".pptx")


def build_folder_output_name(uploaded_files) -> str:
    folder_name = get_common_folder_name(file.name for file in uploaded_files)
    if folder_name:
        return f"{folder_name} RESULT.pptx"
    return build_output_name(uploaded_files)


def build_named_output_name(name: str, extension: str) -> str:
    base_name = name.strip()
    if not base_name:
        base_name = "Flatool"
    if base_name.lower().endswith(extension.lower()):
        return base_name
    if base_name.upper().endswith(" RESULT"):
        return f"{base_name}{extension}"
    return f"{base_name} RESULT{extension}"


def get_detected_folder_name(uploaded_files) -> str:
    return get_common_folder_name(file.name for file in uploaded_files)


def build_folder_batch_output_name(uploaded_files) -> str:
    if len(uploaded_files) == 1 and uploaded_files[0].name.lower().endswith(".zip"):
        return build_result_name(uploaded_files[0].name, ".zip")

    folder_name = get_detected_folder_name(uploaded_files)
    if folder_name:
        return f"{folder_name} RESULT.zip"
    return "Flatool Folder Batch RESULT.zip"


def get_folder_batch_group_names(uploaded_files) -> list[str]:
    groups = group_files_by_child_folder(expand_zip_uploads(uploaded_files))
    return sorted(groups, key=natural_sort_key)


def folder_batch_has_structure(uploaded_files) -> bool:
    if any(file.name.lower().endswith(".zip") for file in uploaded_files):
        return len(get_folder_batch_group_names(uploaded_files)) > 0

    return len(get_folder_batch_group_names(uploaded_files)) > 1


def build_folder_batch_zip(uploaded_files) -> io.BytesIO:
    output = io.BytesIO()
    groups = group_files_by_child_folder(expand_zip_uploads(uploaded_files))

    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as archive:
        for group_name in sorted(groups, key=natural_sort_key):
            pptx_bytes = build_locked_pptx(groups[group_name])
            archive.writestr(f"{group_name} RESULT.pptx", pptx_bytes.getvalue())

    output.seek(0)
    return output


def expand_zip_uploads(uploaded_files) -> list:
    expanded_files = []

    for uploaded_file in uploaded_files:
        if uploaded_file.name.lower().endswith(".zip"):
            with zipfile.ZipFile(io.BytesIO(uploaded_file.getvalue())) as archive:
                for item in archive.infolist():
                    if item.is_dir() or not is_supported_material(item.filename):
                        continue
                    expanded_files.append(MemoryUpload(item.filename, archive.read(item)))
        else:
            expanded_files.append(uploaded_file)

    return expanded_files


def is_supported_material(file_name: str) -> bool:
    return file_name.lower().endswith((".png", ".jpg", ".jpeg", ".pdf"))


def group_files_by_child_folder(uploaded_files) -> dict[str, list]:
    file_paths = [normalize_path_parts(file.name) for file in uploaded_files]
    common_parent = common_parent_parts(file_paths)
    groups: dict[str, list] = {}

    for uploaded_file, parts in zip(uploaded_files, file_paths):
        relative_parts = parts[len(common_parent) :]
        if len(relative_parts) > 1:
            group_name = relative_parts[0]
        elif common_parent:
            group_name = common_parent[-1]
        else:
            group_name = "Flatool"
        groups.setdefault(group_name, []).append(uploaded_file)

    return groups


def common_parent_parts(file_paths: list[list[str]]) -> list[str]:
    if not file_paths:
        return []

    common_parts = file_paths[0][:-1]
    for parts in file_paths[1:]:
        common_parts = common_path_prefix(common_parts, parts[:-1])
    return common_parts


def get_common_folder_name(file_names) -> str:
    path_parts = [
        normalize_path_parts(file_name)
        for file_name in file_names
        if "/" in file_name or "\\" in file_name
    ]

    if not path_parts:
        return ""

    common_parts = path_parts[0][:-1]
    for parts in path_parts[1:]:
        common_parts = common_path_prefix(common_parts, parts[:-1])

    if common_parts:
        return common_parts[-1]
    return path_parts[0][0] if path_parts[0] else ""


def normalize_path_parts(file_name: str) -> list[str]:
    return [part for part in file_name.replace("\\", "/").split("/") if part]


def common_path_prefix(first: list[str], second: list[str]) -> list[str]:
    prefix = []
    for left, right in zip(first, second):
        if left != right:
            break
        prefix.append(left)
    return prefix


def build_result_name(file_name: str, extension: str) -> str:
    base_name = file_name.rsplit(".", 1)[0].strip() or "Flatool"
    return f"{base_name} RESULT{extension}"


def background_xml(relationship_id: str) -> str:
    return f"""<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:bgPr><a:blipFill><a:blip r:embed="{relationship_id}"><a:extLst><a:ext uri="{{28A0092B-C50C-407E-A947-70E740481C1C}}"><a14:useLocalDpi xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main" val="0"/></a:ext></a:extLst></a:blip><a:stretch><a:fillRect/></a:stretch></a:blipFill></p:bgPr></p:bg>"""


def add_background_slide(
    presentation: Presentation,
    image_bytes: bytes,
    width_px: int,
    height_px: int,
    is_first_page: bool,
) -> bool:
    if is_first_page:
        presentation.slide_width = Inches(width_px / 96)
        presentation.slide_height = Inches(height_px / 96)
        is_first_page = False

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _, relationship_id = slide.part.get_or_add_image_part(io.BytesIO(image_bytes))
    slide.element.cSld.insert(0, parse_xml(background_xml(relationship_id)))
    return is_first_page


def render_pdf_pages(pdf_bytes: bytes) -> Iterable[tuple[bytes, int, int]]:
    with fitz.open(stream=pdf_bytes, filetype="pdf") as document:
        for page in document:
            pixmap = page.get_pixmap()
            yield pixmap.tobytes("png"), pixmap.width, pixmap.height


def get_image_size(image_bytes: bytes) -> tuple[int, int]:
    with Image.open(io.BytesIO(image_bytes)) as image:
        return image.size


def build_locked_pptx(uploaded_files) -> io.BytesIO:
    presentation = Presentation()
    is_first_page = True

    for uploaded_file in sorted(uploaded_files, key=lambda file: natural_sort_key(file.name)):
        file_bytes = uploaded_file.getvalue()
        file_name = uploaded_file.name.lower()

        if file_name.endswith(".pdf"):
            for image_bytes, width_px, height_px in render_pdf_pages(file_bytes):
                is_first_page = add_background_slide(
                    presentation,
                    image_bytes,
                    width_px,
                    height_px,
                    is_first_page,
                )
        else:
            width_px, height_px = get_image_size(file_bytes)
            is_first_page = add_background_slide(
                presentation,
                file_bytes,
                width_px,
                height_px,
                is_first_page,
            )

    output = io.BytesIO()
    presentation.save(output)
    output.seek(0)
    return output


def build_pdf_batch_zip(uploaded_files) -> io.BytesIO:
    output = io.BytesIO()

    with zipfile.ZipFile(output, "w", zipfile.ZIP_DEFLATED) as archive:
        for uploaded_file in sorted(uploaded_files, key=lambda file: natural_sort_key(file.name)):
            pptx_bytes = build_locked_pptx([uploaded_file])
            archive.writestr(build_result_name(uploaded_file.name, ".pptx"), pptx_bytes.getvalue())

    output.seek(0)
    return output
