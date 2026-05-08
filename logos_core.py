import hashlib
import io
import re
from typing import Iterable

import fitz
from PIL import Image
from pptx import Presentation
from pptx.oxml import parse_xml
from pptx.util import Inches


VALID_LICENSE_HASHES = {
    hashlib.sha256("LOGOS-BETA-2026".encode()).hexdigest(),
    hashlib.sha256("CCC-FOUNDER-2026".encode()).hexdigest(),
}


def natural_sort_key(value: str) -> list[object]:
    return [
        int(part) if part.isdigit() else part.lower()
        for part in re.split(r"([0-9]+)", value)
    ]


def is_valid_license(license_key: str) -> bool:
    normalized = license_key.strip().upper()
    return hashlib.sha256(normalized.encode()).hexdigest() in VALID_LICENSE_HASHES


def background_xml(relationship_id: str) -> str:
    return f"""<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:bgPr><a:blipFill><a:blip r:embed="{relationship_id}"><a:extLst><a:ext uri="{{28A0092B-C50C-407E-A947-70E740481C1C}}"><a14:useLocalDpi xmlns:a14="http://schemas.microsoft.com/office/drawing/2010/main" val="0"/></a:ext></a:extLst></a:blip><a:stretch><a:fillRect/></a:stretch></a:blipFill></p:bgPr></p:bg>"""


def add_background_slide(
    presentation: Presentation,
    image_bytes: bytes,
    width_px: int,
    height_px: int,
    is_first_page: bool,
) -> bool:
    if is_first_page:
        presentation.slide_width = Inches(width_px / 96)
        presentation.slide_height = Inches(height_px / 96)
        is_first_page = False

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _, relationship_id = slide.part.get_or_add_image_part(io.BytesIO(image_bytes))
    slide.element.cSld.insert(0, parse_xml(background_xml(relationship_id)))
    return is_first_page


def render_pdf_pages(pdf_bytes: bytes) -> Iterable[tuple[bytes, int, int]]:
    with fitz.open(stream=pdf_bytes, filetype="pdf") as document:
        for page in document:
            pixmap = page.get_pixmap()
            yield pixmap.tobytes("png"), pixmap.width, pixmap.height


def get_image_size(image_bytes: bytes) -> tuple[int, int]:
    with Image.open(io.BytesIO(image_bytes)) as image:
        return image.size


def build_locked_pptx(uploaded_files) -> io.BytesIO:
    presentation = Presentation()
    is_first_page = True

    for uploaded_file in sorted(uploaded_files, key=lambda file: natural_sort_key(file.name)):
        file_bytes = uploaded_file.getvalue()
        file_name = uploaded_file.name.lower()

        if file_name.endswith(".pdf"):
            for image_bytes, width_px, height_px in render_pdf_pages(file_bytes):
                is_first_page = add_background_slide(
                    presentation,
                    image_bytes,
                    width_px,
                    height_px,
                    is_first_page,
                )
        else:
            width_px, height_px = get_image_size(file_bytes)
            is_first_page = add_background_slide(
                presentation,
                file_bytes,
                width_px,
                height_px,
                is_first_page,
            )

    output = io.BytesIO()
    presentation.save(output)
    output.seek(0)
    return output

